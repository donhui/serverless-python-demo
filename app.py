from pptx import Presentation
import urllib.request

def ppt_to_text(path):
    prs = Presentation(r''+path);
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                # print(''.join(run.text for run in paragraph.runs))
                return ''.join(run.text for run in paragraph.runs)

def download(url, path, name):
    urllib.request.urlretrieve(url, '{0}{1}'.format(path, name))
    return path+name